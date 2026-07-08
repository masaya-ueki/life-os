"""deckgen のスモークテスト。

- 実 deck(claude-code-security)からスライド数一致・表シェイプ存在を確認
- 合成 outline で全 expression（chart/emphasis/structure/flow/comparison）を描画し、
  チャートが生成され、画像が一切含まれない（=編集可能ネイティブ）ことを確認
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Pt

from deckgen import layout
from deckgen.builder import build_presentation
from deckgen.loader import iter_slides, load_outline
from deckgen.theme import get_theme


def _all_shapes(prs):
    for s in prs.slides:
        yield from s.shapes


def test_real_deck_slide_count_and_table():
    outline = load_outline("claude-code-security")
    expected = sum(1 for _ in iter_slides(outline))
    prs, warnings = build_presentation(outline)
    assert len(prs.slides) == expected
    # comparison table を含むスライドがある
    assert any(sh.has_table for sh in _all_shapes(prs))
    # 既知 expression のみなので未知警告は無い
    assert warnings == []


def test_no_images_means_editable():
    outline = load_outline("claude-code-security")
    prs, _ = build_presentation(outline)
    pics = [
        sh for sh in _all_shapes(prs)
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert pics == []


SYNTH = {
    "deck": {"title": "テスト", "subtitle": "sub", "theme": "dark"},
    "chapters": [
        {
            "chapter": "c",
            "slides": [
                {"title": "表紙", "summary": "s", "content": ["a"],
                 "expression": "title"},
                {"title": "棒グラフ", "summary": "推移", "expression": "chart",
                 "data": {"type": "bar", "unit": "%",
                          "series": [{"label": "2024", "value": 30},
                                     {"label": "2025", "value": 65}]}},
                {"title": "強調", "summary": "msg", "expression": "emphasis",
                 "data": {"mode": "big-number", "value": "-80", "unit": "%",
                          "label": "件数"}},
                {"title": "マトリクス", "summary": "m", "expression": "structure",
                 "data": {"type": "matrix-2x2", "axis_x": "x", "axis_y": "y",
                          "quadrants": ["A", "B", "C", "D"]}},
                {"title": "フロー", "summary": "f", "expression": "flow",
                 "data": {"type": "steps", "orientation": "horizontal",
                          "steps": [{"label": "1", "desc": "d"},
                                    {"label": "2", "desc": "d"}]}},
                {"title": "未知", "summary": "u", "content": ["x", "y"],
                 "expression": "nonexistent"},
            ],
        }
    ],
}


def test_synthetic_all_expressions_and_chart():
    prs, warnings = build_presentation(SYNTH)
    assert len(prs.slides) == 6
    # ネイティブチャートが1つ生成される
    assert sum(1 for sh in _all_shapes(prs) if sh.has_chart) == 1
    # 画像は無い
    assert not any(
        sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in _all_shapes(prs)
    )
    # 未知 expression は警告される
    assert any("nonexistent" in w for w in warnings)


def test_titles_present():
    prs, _ = build_presentation(SYNTH)
    texts = " ".join(
        sh.text_frame.text for sh in _all_shapes(prs) if sh.has_text_frame
    )
    assert "棒グラフ" in texts
    assert "マトリクス" in texts


# --- chart 実運用検証 + data 契約の回帰テスト (#33) ---

CHART_DECK = "deckgen-chart-demo"


def _count_charts(prs):
    return sum(1 for sh in _all_shapes(prs) if sh.has_chart)


def _chart_outline(ctype, series, *, unit="%", content=None):
    """chart 1枚だけの最小 outline を組み立てる（契約テスト用）。"""
    slide = {"title": "chart", "expression": "chart",
             "data": {"type": ctype, "unit": unit, "series": series}}
    if content is not None:
        slide["content"] = content
    return {
        "deck": {"title": "t", "theme": "default"},
        "chapters": [{"chapter": "c", "slides": [slide]}],
    }


def test_chart_demo_deck_renders_native_charts():
    """chart サンプル deck(bar/line/pie/stacked)が実運用でネイティブ生成される。"""
    outline = load_outline(CHART_DECK)
    expected = sum(1 for _ in iter_slides(outline))
    prs, warnings = build_presentation(outline)
    assert len(prs.slides) == expected
    # 表紙1 + チャート4枚（bar/line/pie/stacked）
    assert _count_charts(prs) == 4
    assert warnings == []
    # 画像なし＝編集可能ネイティブ
    assert not any(
        sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in _all_shapes(prs)
    )


def test_chart_contract_single_series_builds_chart():
    """単系列(bar/pie): label/value 契約でネイティブチャートになる。"""
    series = [{"label": "A", "value": 1}, {"label": "B", "value": 2}]
    for ctype in ("bar", "pie"):
        prs, _ = build_presentation(_chart_outline(ctype, series))
        assert _count_charts(prs) == 1, ctype


def test_chart_contract_multi_series_builds_chart():
    """複数系列(line/stacked): name/points(x,y)契約でネイティブチャートになる。"""
    series = [{"name": "S", "points": [{"x": "1", "y": 1}, {"x": "2", "y": 3}]}]
    for ctype in ("line", "stacked"):
        prs, _ = build_presentation(_chart_outline(ctype, series))
        assert _count_charts(prs) == 1, ctype


def test_chart_empty_series_falls_back_to_bullet():
    """series 空 → チャートを作らず本文へフォールバック（契約破綻の検知）。"""
    prs, _ = build_presentation(
        _chart_outline("bar", [], content=["フォールバック本文"])
    )
    assert _count_charts(prs) == 0
    texts = " ".join(
        sh.text_frame.text for sh in _all_shapes(prs) if sh.has_text_frame
    )
    assert "フォールバック本文" in texts


def test_chart_multi_series_without_points_falls_back():
    """複数系列で points が無い → カテゴリ不能でフォールバック（チャート無し）。"""
    prs, _ = build_presentation(_chart_outline("line", [{"name": "S"}]))
    assert _count_charts(prs) == 0


def test_html_pptx_slide_count_parity():
    """同一 outline.yml の論理スライド数と生成 pptx の枚数が一致する。

    HTML レンダラは決定的コードを持たない（エージェント生成）ため、HTML/pptx の
    整合は「outline.yml という単一の真実から導かれるスライド数の一致」で担保する。
    outline 側の章/枚数変更や、ビルダのスライド生成漏れを検知する。
    """
    for slug in ("claude-code-security", CHART_DECK):
        outline = load_outline(slug)
        prs, _ = build_presentation(outline)
        assert len(prs.slides) == sum(1 for _ in iter_slides(outline)), slug


# --- カードスタイル統一の契約テスト (#109) ---


def _fill_rgb(shape):
    """ソリッド塗りの RGB を返す。ソリッド塗りでなければ None（例外を握りつぶさない）。"""
    try:
        return shape.fill.fore_color.rgb
    except (AttributeError, TypeError, ValueError):
        return None


FLOW_OUTLINE = {
    "deck": {"title": "t", "theme": "default"},
    "chapters": [
        {
            "chapter": "c",
            "slides": [
                {"title": "flow", "summary": "s", "expression": "flow",
                 "data": {"type": "steps", "orientation": "horizontal",
                          "steps": [{"label": "1", "desc": "d1"},
                                    {"label": "2", "desc": "d2"}]}},
            ],
        }
    ],
}


def test_flow_step_cards_use_standard_card_border():
    """flow ステップの枠線が accent 強調ではなく標準カード（line 色・Pt(1.0)）に統一されている。"""
    theme = get_theme("default")
    prs, _ = build_presentation(FLOW_OUTLINE)
    card_shapes = [
        sh for sh in _all_shapes(prs)
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and sh.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
        and _fill_rgb(sh) == RGBColor.from_string(theme["card"])
    ]
    assert card_shapes
    for sh in card_shapes:
        assert sh.line.color.rgb == RGBColor.from_string(theme["line"])
        assert sh.line.width == Pt(1.0)


MATRIX_OUTLINE = {
    "deck": {"title": "t", "theme": "default"},
    "chapters": [
        {
            "chapter": "c",
            "slides": [
                {"title": "matrix", "summary": "s", "expression": "structure",
                 "data": {"type": "matrix-2x2", "axis_x": "x", "axis_y": "y",
                          "quadrants": ["A", "B", "C", "D"]}},
            ],
        }
    ],
}


def test_matrix_quadrants_are_rounded_cards():
    """matrix-2x2 の象限セルが（矩形ではなく）角丸カードで描かれている。"""
    prs, _ = build_presentation(MATRIX_OUTLINE)
    quadrant_shapes = [
        sh for sh in _all_shapes(prs)
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and sh.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
        and sh.has_text_frame and sh.text_frame.text in ("A", "B", "C", "D")
    ]
    assert len(quadrant_shapes) == 4
    for sh in quadrant_shapes:
        # adjustments はネイティブ角丸(roundRect)のみが持つ調整ハンドル
        assert sh.adjustments[0] > 0


KPI_OUTLINE = {
    "deck": {"title": "t", "theme": "default"},
    "chapters": [
        {
            "chapter": "c",
            "slides": [
                {"title": "kpi", "summary": "s", "expression": "emphasis",
                 "data": {"mode": "kpi", "cards": [
                     {"num": "120", "delta": "-12%", "label": "件数"},
                     {"num": "98", "delta": "+5%", "label": "率"},
                 ]}},
            ],
        }
    ],
}


def test_kpi_delta_uses_semantic_colors():
    """KPI の増減表示がマイナスは bad、プラスは good の意味色になる。"""
    theme = get_theme("default")
    prs, _ = build_presentation(KPI_OUTLINE)
    runs = [
        r
        for sh in _all_shapes(prs) if sh.has_text_frame
        for p in sh.text_frame.paragraphs
        for r in p.runs
    ]
    bad_run = next(r for r in runs if r.text == " -12%")
    good_run = next(r for r in runs if r.text == " +5%")
    assert bad_run.font.color.rgb == RGBColor.from_string(theme["bad"])
    assert good_run.font.color.rgb == RGBColor.from_string(theme["good"])


TITLE_OUTLINE = {
    "deck": {"title": "タイトル", "subtitle": "サブ", "theme": "default", "date": "2024-01-01"},
    "chapters": [
        {
            "chapter": "c",
            "slides": [
                {"title": "表紙", "summary": "s", "content": ["a", "b"],
                 "expression": "title"},
            ],
        }
    ],
}


def test_title_cover_shapes_align_to_content_left():
    """表紙のテキストを持つ図形はすべてコンテンツスライドと同じ CONTENT_LEFT に左揃えされる。"""
    prs, _ = build_presentation(TITLE_OUTLINE)
    text_shapes = [
        sh for sh in _all_shapes(prs)
        if sh.has_text_frame and sh.text_frame.text.strip()
    ]
    assert text_shapes
    for sh in text_shapes:
        assert sh.left == int(layout.CONTENT_LEFT)
