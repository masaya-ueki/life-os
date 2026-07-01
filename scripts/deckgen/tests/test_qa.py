"""deckgen.qa のテスト。

正常系: 現存 deck が QA を pass すること
異常系: 合成データで各 check が正しく発火すること
"""

from __future__ import annotations

import tempfile
from pathlib import Path

import pytest

from deckgen.qa import (
    Finding,
    check_html,
    check_outline,
    check_pptx,
)


# ---------------------------------------------------------------------------
# Outline チェック
# ---------------------------------------------------------------------------

def _make_outline(slides: list[dict]) -> dict:
    return {
        "deck": {"title": "テスト", "theme": "default"},
        "chapters": [{"chapter": "c", "slides": slides}],
    }


def test_outline_clean_passes():
    outline = _make_outline([
        {"title": "タイトル", "summary": "s", "content": ["a", "b"], "expression": "bullet"},
    ])
    assert check_outline(outline, "test/outline.yml") == []


def test_outline_empty_title_flagged():
    outline = _make_outline([{"title": "", "expression": "bullet"}])
    findings = check_outline(outline, "test/outline.yml")
    assert any(f.check == "QA-OUTLINE-EMPTY" and f.severity == "error" for f in findings)


def test_outline_whitespace_title_flagged():
    outline = _make_outline([{"title": "   ", "expression": "bullet"}])
    findings = check_outline(outline, "test/outline.yml")
    assert any(f.check == "QA-OUTLINE-EMPTY" for f in findings)


def test_outline_placeholder_in_title_flagged():
    outline = _make_outline([{"title": "{{タイトルを入力}}", "expression": "bullet"}])
    findings = check_outline(outline, "test/outline.yml")
    assert any(f.check == "QA-OUTLINE-PH" and f.severity == "warning" for f in findings)


def test_outline_placeholder_in_content_flagged():
    outline = _make_outline([
        {"title": "t", "content": ["正常", "{{内容を入力}}"], "expression": "bullet"},
    ])
    findings = check_outline(outline, "test/outline.yml")
    assert any(f.check == "QA-OUTLINE-PH" for f in findings)


def test_outline_excess_bullets_flagged():
    """content が 8 件以上で QA-OUTLINE-EXCESS-BULLETS が発火する。"""
    content = [f"item{i}" for i in range(8)]  # 8 件（MAX_BULLETS=7 超え）
    outline = _make_outline([{"title": "t", "content": content, "expression": "bullet"}])
    findings = check_outline(outline, "test/outline.yml")
    assert any(f.check == "QA-OUTLINE-EXCESS-BULLETS" and f.severity == "warning" for f in findings)


def test_outline_exactly_max_bullets_passes():
    """content が MAX_BULLETS(7) 件ちょうどはフラグなし。"""
    content = [f"item{i}" for i in range(7)]
    outline = _make_outline([{"title": "t", "content": content, "expression": "bullet"}])
    findings = check_outline(outline, "test/outline.yml")
    assert not any(f.check == "QA-OUTLINE-EXCESS-BULLETS" for f in findings)


def test_real_deck_outline_passes():
    """claude-code-security の outline.yml は品質問題なし。"""
    from deckgen.loader import load_outline
    outline = load_outline("claude-code-security")
    findings = check_outline(outline, "claude-code-security/outline.yml")
    errors = [f for f in findings if f.severity == "error"]
    assert errors == [], errors


# ---------------------------------------------------------------------------
# HTML チェック
# ---------------------------------------------------------------------------

_HTML_TEMPLATE = """\
<!DOCTYPE html>
<html lang="ja">
<head><meta charset="UTF-8"><title>test</title><style>body{{}}</style></head>
<body>
{slides}
<script>var x=1;</script>
</body>
</html>
"""

_SLIDE_TEMPLATE = '<div class="slide" data-variant="{variant}">{content}</div>'


def _make_html(slides_html: str) -> str:
    return _HTML_TEMPLATE.format(slides=slides_html)


def test_html_clean_passes():
    slides = _SLIDE_TEMPLATE.format(variant="bullet", content="<h2>タイトル</h2><p>内容</p>")
    findings = check_html_from_str(_make_html(slides), "test/index.html")
    assert findings == []


def test_html_external_script_flagged():
    html = '<!DOCTYPE html><html><head><script src="https://cdn.example.com/app.js"></script></head><body></body></html>'
    findings = check_html_from_str(html, "test/index.html")
    assert any(f.check == "QA-HTML-EXTERNAL" and f.severity == "error" for f in findings)


def test_html_external_link_flagged():
    html = '<!DOCTYPE html><html><head><link href="https://fonts.googleapis.com/css" rel="stylesheet"></head><body></body></html>'
    findings = check_html_from_str(html, "test/index.html")
    assert any(f.check == "QA-HTML-EXTERNAL" for f in findings)


def test_html_external_img_src_in_slide_flagged():
    content = '<img src="https://example.com/img.png" alt="x">'
    slides = _SLIDE_TEMPLATE.format(variant="bullet", content=content)
    findings = check_html_from_str(_make_html(slides), "test/index.html")
    assert any(f.check == "QA-HTML-EXTERNAL" for f in findings)


def test_html_placeholder_in_slide_flagged():
    content = "<h2>{{スライドタイトル}}</h2>"
    slides = _SLIDE_TEMPLATE.format(variant="bullet", content=content)
    findings = check_html_from_str(_make_html(slides), "test/index.html")
    assert any(f.check == "QA-HTML-EMPTY-PH" and f.severity == "warning" for f in findings)


def test_html_overflow_in_slide_flagged():
    """601 文字超のテキストを含むスライドで QA-HTML-OVERFLOW が発火する。"""
    long_text = "あ" * 601
    content = f"<p>{long_text}</p>"
    slides = _SLIDE_TEMPLATE.format(variant="bullet", content=content)
    findings = check_html_from_str(_make_html(slides), "test/index.html")
    assert any(f.check == "QA-HTML-OVERFLOW" and f.severity == "warning" for f in findings)


def test_html_overflow_exactly_600_passes():
    """ちょうど 600 文字はフラグなし。"""
    text = "あ" * 600
    content = f"<p>{text}</p>"
    slides = _SLIDE_TEMPLATE.format(variant="bullet", content=content)
    findings = check_html_from_str(_make_html(slides), "test/index.html")
    assert not any(f.check == "QA-HTML-OVERFLOW" for f in findings)


def test_real_deck_html_no_external_deps():
    """claude-code-security/index.html に外部 URL がない。"""
    from deckgen.loader import DECKS_DIR
    html_path = DECKS_DIR / "claude-code-security" / "index.html"
    if not html_path.exists():
        pytest.skip("index.html が生成されていない")
    findings = check_html(html_path)
    external = [f for f in findings if f.check == "QA-HTML-EXTERNAL"]
    assert external == [], external


# ---------------------------------------------------------------------------
# PPTX チェック
# ---------------------------------------------------------------------------

def _make_minimal_pptx(out: Path) -> None:
    """python-pptx で最小の pptx（テキストフレームなし）を生成する。"""
    from pptx import Presentation
    prs = Presentation()
    layout = prs.slide_layouts[6]  # blank
    prs.slides.add_slide(layout)
    prs.save(str(out))


def _make_pptx_with_text(out: Path, text: str) -> None:
    """指定テキストを持つ pptx を生成する。"""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    txBox.text_frame.text = text
    prs.save(str(out))


def _make_pptx_with_empty_frame(out: Path) -> None:
    """空テキストフレームを含む pptx を生成する。"""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    txBox.text_frame.text = ""
    prs.save(str(out))


def test_pptx_clean_passes():
    with tempfile.TemporaryDirectory() as tmpdir:
        out = Path(tmpdir) / "test.pptx"
        _make_pptx_with_text(out, "正常なテキスト")
        findings = check_pptx(out)
        errors = [f for f in findings if f.severity == "error"]
        assert errors == []


def test_pptx_empty_text_frame_flagged():
    """空テキストフレームで QA-PPTX-EMPTY が発火する。"""
    with tempfile.TemporaryDirectory() as tmpdir:
        out = Path(tmpdir) / "test.pptx"
        _make_pptx_with_empty_frame(out)
        findings = check_pptx(out)
        assert any(f.check == "QA-PPTX-EMPTY" and f.severity == "warning" for f in findings)


def test_pptx_placeholder_flagged():
    """プレースホルダパターンで QA-PPTX-PH が発火する。"""
    with tempfile.TemporaryDirectory() as tmpdir:
        out = Path(tmpdir) / "test.pptx"
        _make_pptx_with_text(out, "{{ここに内容を入力}}")
        findings = check_pptx(out)
        assert any(f.check == "QA-PPTX-PH" and f.severity == "warning" for f in findings)


def test_real_deck_pptx_no_errors():
    """claude-code-security から生成した pptx に error がない。"""
    from deckgen.builder import build_to_file
    from deckgen.loader import load_outline
    outline = load_outline("claude-code-security")
    with tempfile.TemporaryDirectory() as tmpdir:
        out = Path(tmpdir) / "claude-code-security.pptx"
        build_to_file(outline, out)
        findings = check_pptx(out)
        errors = [f for f in findings if f.severity == "error"]
        assert errors == [], errors


# ---------------------------------------------------------------------------
# 視覚品質チェック（コントラスト・フォント・色トークン・境界・重なり）
# ---------------------------------------------------------------------------

def _make_pptx_runs(out: Path, runs: list[tuple], *, box=(1, 1, 8, 5)) -> None:
    """1 テキストボックスに複数 run を持つ pptx を生成する。
    runs: (text, size_pt|None, color_hex|None) のリスト。"""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(
        Inches(box[0]), Inches(box[1]), Inches(box[2]), Inches(box[3]))
    p = tb.text_frame.paragraphs[0]
    for text, size, color in runs:
        r = p.add_run()
        r.text = text
        if size is not None:
            r.font.size = Pt(size)
        if color is not None:
            r.font.color.rgb = RGBColor.from_string(color)
    prs.save(str(out))


def test_contrast_ratio_extremes():
    from deckgen.qa import contrast_ratio
    assert round(contrast_ratio("#000000", "#ffffff"), 1) == 21.0
    assert round(contrast_ratio("#ffffff", "#ffffff"), 1) == 1.0


def test_theme_contrast_flags_low_contrast():
    """good #16A34A（白背景 3.3:1）等の低コントラストを検出する。"""
    from deckgen.qa import check_theme_contrast
    theme = {"fg": "1A1A2E", "bg": "FFFFFF", "card": "FFFFFF",
             "muted": "6B7280", "good": "16A34A", "bad": "DC2626",
             "accent": "2563EB", "on_accent": "FFFFFF"}
    findings = check_theme_contrast(theme, "t/theme")
    assert any(f.check == "QA-THEME-CONTRAST" for f in findings)
    assert any("good on bg" in f.message for f in findings)


def test_theme_contrast_high_contrast_passes():
    from deckgen.qa import check_theme_contrast
    theme = {"fg": "111111", "bg": "FFFFFF", "card": "FFFFFF",
             "muted": "4B5563", "good": "166534", "bad": "B91C1C",
             "accent": "1D4ED8", "on_accent": "FFFFFF"}
    findings = check_theme_contrast(theme, "t/theme")
    assert findings == [], [f.message for f in findings]


def test_pptx_minfont_flagged():
    with tempfile.TemporaryDirectory() as d:
        out = Path(d) / "t.pptx"
        _make_pptx_runs(out, [("小さすぎる文字", 10, None)])
        findings = check_pptx(out)
        assert any(f.check == "QA-PPTX-MINFONT" for f in findings)


def test_pptx_typescale_flagged():
    with tempfile.TemporaryDirectory() as d:
        out = Path(d) / "t.pptx"
        _make_pptx_runs(out, [("型スケール外", 17, None)])  # 17pt は TYPE_SCALE 外
        findings = check_pptx(out)
        assert any(f.check == "QA-PPTX-TYPESCALE" for f in findings)


def test_pptx_typescale_in_scale_passes():
    with tempfile.TemporaryDirectory() as d:
        out = Path(d) / "t.pptx"
        _make_pptx_runs(out, [("本文", 18, None)])  # FONT_BODY
        findings = check_pptx(out)
        assert not any(f.check in ("QA-PPTX-TYPESCALE", "QA-PPTX-MINFONT")
                       for f in findings)


def test_pptx_colortoken_flagged_with_theme():
    theme = {"fg": "1A1A2E", "bg": "FFFFFF", "accent": "2563EB"}
    with tempfile.TemporaryDirectory() as d:
        out = Path(d) / "t.pptx"
        _make_pptx_runs(out, [("テーマ外の色", 18, "123456")])
        findings = check_pptx(out, theme=theme)
        assert any(f.check == "QA-PPTX-COLORTOKEN" for f in findings)


def test_pptx_colortoken_skipped_without_theme():
    with tempfile.TemporaryDirectory() as d:
        out = Path(d) / "t.pptx"
        _make_pptx_runs(out, [("色", 18, "123456")])
        findings = check_pptx(out)  # theme=None
        assert not any(f.check == "QA-PPTX-COLORTOKEN" for f in findings)


def test_pptx_bounds_flagged():
    with tempfile.TemporaryDirectory() as d:
        out = Path(d) / "t.pptx"
        _make_pptx_runs(out, [("はみ出し", 18, None)], box=(14, 1, 4, 2))
        findings = check_pptx(out)
        assert any(f.check == "QA-PPTX-BOUNDS" for f in findings)


def test_pptx_sizekinds_flagged():
    with tempfile.TemporaryDirectory() as d:
        out = Path(d) / "t.pptx"
        runs = [("a", 14, None), ("b", 16, None), ("c", 18, None),
                ("d", 22, None), ("e", 36, None), ("f", 40, None)]  # 6 種 > MAX_SIZE_KINDS(5)
        _make_pptx_runs(out, runs)
        findings = check_pptx(out)
        assert any(f.check == "QA-PPTX-SIZEKINDS" for f in findings)


def test_pptx_overlap_flagged():
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    with tempfile.TemporaryDirectory() as d:
        out = Path(d) / "t.pptx"
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        for x in (1.0, 2.0):  # 大きく重なる 2 枚の矩形
            sp = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(1), Inches(4), Inches(3))
            sp.fill.solid()
            sp.fill.fore_color.rgb = RGBColor.from_string("2563EB")
        prs.save(str(out))
        findings = check_pptx(out)
        assert any(f.check == "QA-PPTX-OVERLAP" for f in findings)


# ---------------------------------------------------------------------------
# ヘルパ（テスト内のみ使用）
# ---------------------------------------------------------------------------

def check_html_from_str(html: str, label: str) -> list[Finding]:
    """HTML 文字列を直接検査するテスト用ヘルパ。"""
    from deckgen.qa import _HtmlQaParser
    parser = _HtmlQaParser(label)
    parser.feed(html)
    return parser.findings
