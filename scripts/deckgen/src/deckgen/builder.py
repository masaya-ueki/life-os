"""outline.yml(dict) → python-pptx Presentation の組み立て。

- deck.theme で配色を決める（テンプレ未指定時は背景を塗る）
- 各スライド共通: タイトル(h2) + リード(summary) のヘッダ
- expression に応じて本文を描画（未知/欠落は bullet にフォールバック）
- title は全面の表紙レイアウト
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from deckgen import layout
from deckgen.expressions import get_renderer
from deckgen.expressions import bullet, title
from deckgen.loader import iter_slides
from deckgen.theme import get_theme


def build_presentation(outline: dict, template_path: str | None = None):
    """outline(dict) から Presentation を構築して返す。"""
    prs = Presentation(template_path) if template_path else Presentation()
    prs.slide_width = layout.SLIDE_W
    prs.slide_height = layout.SLIDE_H

    deck = outline.get("deck", {})
    theme = get_theme(deck.get("theme"))
    use_bg = template_path is None  # テンプレ使用時はマスター背景に任せる

    warnings: list[str] = []

    for chapter, slide in iter_slides(outline):
        pslide = layout.add_slide(prs)
        if use_bg:
            layout.fill_background(pslide, theme["bg"])

        expression = slide.get("expression")

        if expression == "title":
            title.render_cover(pslide, theme, deck, slide)
            continue

        region = layout.add_header(
            pslide, theme, slide.get("title", ""), slide.get("summary", "")
        )

        renderer = get_renderer(expression)
        if renderer is None:
            if expression not in (None, "", "bullet"):
                warnings.append(
                    f"未知の expression '{expression}'（{slide.get('title')}）→ bullet で描画"
                )
            renderer = bullet.render
        renderer(pslide, theme, slide, region)

    return prs, warnings


def build_to_file(outline: dict, out_path: str | Path, template_path: str | None = None):
    """構築して .pptx を書き出し、(出力パス, スライド数, 警告) を返す。"""
    prs, warnings = build_presentation(outline, template_path)
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out, len(prs.slides), warnings
