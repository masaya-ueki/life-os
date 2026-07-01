"""生成スライドの QA チェッカー（post-generation fitness function）。

outline.yml / index.html / .pptx の 3 つの生成物を検査し、
明らかな品質問題（テキスト溢れ・空プレースホルダ・外部依存・1スライド1メッセージ逸脱）
を自動検知・報告する。check_structure.py と同じ出力インタフェースを採用する。

使い方:
    uv run --project scripts/deckgen -m deckgen.qa <slug>        # 特定 deck を検査
    uv run --project scripts/deckgen -m deckgen.qa --all          # 全 deck を検査
    uv run --project scripts/deckgen -m deckgen.qa <slug> --json  # JSON 出力
    uv run --project scripts/deckgen -m deckgen.qa <slug> --strict # warning でも exit 1

終了コード: error が 1 件以上なら 1。--strict なら warning でも 1。なければ 0。
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import asdict, dataclass
from html.parser import HTMLParser
from pathlib import Path

from deckgen.loader import DECKS_DIR, load_outline, resolve_outline_path

# --- 定数 ---

PLACEHOLDER_RE = re.compile(r"\{\{[^}]+\}\}")
MAX_BULLETS = 7         # これを超えると 1スライド1メッセージ逸脱と判定
OVERFLOW_CHARS_HTML = 600  # 1スライドのテキスト上限（HTML）
OVERFLOW_CHARS_PER_CM2 = 6  # pptx: 1cm² あたりの推定最大文字数
OVERFLOW_MIN_CHARS = 100    # pptx: 最低閾値

# --- 視覚品質チェックの閾値（pptx） ---
MIN_FONT_PT = 14            # これ未満は可読性の下限割れ（型スケール FONT_CAPTION と一致）
MAX_SIZE_KINDS = 5          # 1 スライドの異なるフォントサイズ種の上限（見出し2 + 多層図解の3層を許容。6 種以上で警告）
CONTRAST_AA_NORMAL = 4.5    # WCAG AA（通常テキスト）
CONTRAST_AA_LARGE = 3.0     # WCAG AA（大きいテキスト）
OVERLAP_AREA_FRAC = 0.4     # 塗り形状の重なりが小さい方の面積のこの割合超で警告
BOUNDS_EPS_IN = 0.03        # 境界超過判定の許容（インチ、丸め誤差の吸収）
THIN_SHAPE_IN = 0.2         # これ未満の辺を持つ塗り形状は装飾（バー/罫線）として重なり判定から除外

# void 要素は handle_endtag が呼ばれないため depth カウントから除外
_VOID_ELEMENTS = frozenset(
    ["area", "base", "br", "col", "embed", "hr", "img", "input",
     "link", "meta", "param", "source", "track", "wbr"]
)


@dataclass
class Finding:
    check: str
    severity: str  # "error" | "warning"
    path: str
    message: str


# ---------------------------------------------------------------------------
# Outline チェック
# ---------------------------------------------------------------------------

def check_outline(outline: dict, label: str) -> list[Finding]:
    """outline.yml を検査し、論理的な品質問題を返す。"""
    findings: list[Finding] = []
    for ci, ch in enumerate(outline.get("chapters", [])):
        for si, slide in enumerate(ch.get("slides", [])):
            slide_label = f"{label}:ch{ci + 1}/slide{si + 1}"

            # QA-OUTLINE-EMPTY: title が空
            title = str(slide.get("title", ""))
            if not title.strip():
                findings.append(Finding(
                    "QA-OUTLINE-EMPTY", "error", slide_label, "title が空"))

            # QA-OUTLINE-PH: プレースホルダ残存
            all_text = " ".join([
                title,
                str(slide.get("summary", "")),
                str(slide.get("message", "")),
                " ".join(str(c) for c in slide.get("content", [])),
            ])
            if PLACEHOLDER_RE.search(all_text):
                findings.append(Finding(
                    "QA-OUTLINE-PH", "warning", slide_label, "プレースホルダ残存"))

            # QA-OUTLINE-EXCESS-BULLETS: 1スライド1メッセージ逸脱
            content = slide.get("content", [])
            if isinstance(content, list) and len(content) > MAX_BULLETS:
                findings.append(Finding(
                    "QA-OUTLINE-EXCESS-BULLETS", "warning", slide_label,
                    f"content が多い ({len(content)} 件 > {MAX_BULLETS})"))

    return findings


# ---------------------------------------------------------------------------
# HTML チェック
# ---------------------------------------------------------------------------

class _HtmlQaParser(HTMLParser):
    """`.slide` 要素を追跡して品質問題を検出する HTMLParser。"""

    def __init__(self, label_prefix: str) -> None:
        super().__init__()
        self._prefix = label_prefix
        self.findings: list[Finding] = []
        self._depth = 0
        self._slide_count = 0
        self._in_slide = False
        self._slide_open_depth = -1
        self._text_buf: list[str] = []

    def handle_starttag(self, tag: str, attrs: list) -> None:
        attrs_dict = dict(attrs)

        # QA-HTML-EXTERNAL: src / href に外部 URL
        for attr in ("src", "href"):
            val = attrs_dict.get(attr, "")
            if val.startswith(("http://", "https://")):
                self.findings.append(Finding(
                    "QA-HTML-EXTERNAL", "error",
                    f"{self._prefix}:{tag}",
                    f"外部 URL ({attr}): {val[:80]}"))

        # QA-HTML-EXTERNAL: style 属性の url() に外部 URL
        for m in re.finditer(r"url\(['\"]?(https?://[^\s'\"\)\]]+)", attrs_dict.get("style", "")):
            self.findings.append(Finding(
                "QA-HTML-EXTERNAL", "error",
                f"{self._prefix}:{tag}",
                f"外部 URL (style): {m.group(1)[:80]}"))

        if tag.lower() in _VOID_ELEMENTS:
            return

        self._depth += 1

        # slide div の開始を検出
        if tag.lower() == "div" and not self._in_slide:
            classes = attrs_dict.get("class", "").split()
            if "slide" in classes:
                self._in_slide = True
                self._slide_open_depth = self._depth
                self._slide_count += 1
                self._text_buf = []

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() in _VOID_ELEMENTS:
            return
        if self._in_slide and tag.lower() == "div" and self._depth == self._slide_open_depth:
            self._flush_slide()
            self._in_slide = False
        self._depth -= 1

    def handle_startendtag(self, tag: str, attrs: list) -> None:
        # XHTML スタイルの自己閉鎖タグ（<div/> など）への対応
        self.handle_starttag(tag, attrs)
        if tag.lower() not in _VOID_ELEMENTS:
            self.handle_endtag(tag)

    def handle_data(self, data: str) -> None:
        if self._in_slide:
            self._text_buf.append(data)

    def _flush_slide(self) -> None:
        text = "".join(self._text_buf)
        label = f"{self._prefix}:slide-{self._slide_count}"

        # QA-HTML-EMPTY-PH: プレースホルダ残存
        if PLACEHOLDER_RE.search(text):
            self.findings.append(Finding(
                "QA-HTML-EMPTY-PH", "warning", label, "プレースホルダ残存"))

        # QA-HTML-OVERFLOW: テキスト量超過
        clean = re.sub(r"\s+", " ", text).strip()
        if len(clean) > OVERFLOW_CHARS_HTML:
            self.findings.append(Finding(
                "QA-HTML-OVERFLOW", "warning", label,
                f"テキスト量超過 ({len(clean)} 文字 > {OVERFLOW_CHARS_HTML})"))


def check_html(html_path: Path, label_prefix: str | None = None) -> list[Finding]:
    """index.html を検査し、レンダリング品質問題を返す。"""
    prefix = label_prefix or html_path.as_posix()
    parser = _HtmlQaParser(prefix)
    parser.feed(html_path.read_text(encoding="utf-8"))
    return parser.findings


# ---------------------------------------------------------------------------
# 視覚品質ヘルパ（WCAG コントラスト・図形）
# ---------------------------------------------------------------------------

def _rel_luminance(hex_color: str) -> float:
    """sRGB hex の WCAG 相対輝度（0.0–1.0）を返す。"""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return 0.0

    def _chan(c: int) -> float:
        x = c / 255.0
        return x / 12.92 if x <= 0.03928 else ((x + 0.055) / 1.055) ** 2.4

    r, g, b = (int(h[i:i + 2], 16) for i in (0, 2, 4))
    return 0.2126 * _chan(r) + 0.7152 * _chan(g) + 0.0722 * _chan(b)


def contrast_ratio(hex1: str, hex2: str) -> float:
    """2 色の WCAG コントラスト比（1.0–21.0）を返す。"""
    l1, l2 = _rel_luminance(hex1), _rel_luminance(hex2)
    hi, lo = max(l1, l2), min(l1, l2)
    return (hi + 0.05) / (lo + 0.05)


# レンダラが実際に生成する「意味色 × 面」の組。逸脱は再利用色の a11y 問題。
_CONTRAST_PAIRS = (
    ("fg", "bg"),          # 本文
    ("fg", "card"),
    ("muted", "bg"),       # 補足・注記
    ("muted", "card"),
    ("good", "bg"),        # メリット（pros-cons / KPI delta）
    ("good", "card"),
    ("bad", "bg"),         # デメリット
    ("bad", "card"),
    ("accent", "bg"),      # 見出し・KPI 数値
    ("on_accent", "accent"),  # アクセント面上のラベル
)


def check_theme_contrast(theme: dict, label: str) -> list[Finding]:
    """テーマの意味色ペアが WCAG AA(4.5:1) を満たすか検査する（deck 非依存）。"""
    findings: list[Finding] = []
    for text_tok, surf_tok in _CONTRAST_PAIRS:
        if text_tok not in theme or surf_tok not in theme:
            continue
        ratio = contrast_ratio(theme[text_tok], theme[surf_tok])
        if ratio < CONTRAST_AA_NORMAL:
            findings.append(Finding(
                "QA-THEME-CONTRAST", "warning", label,
                f"{text_tok} on {surf_tok} = {ratio:.2f}:1 < {CONTRAST_AA_NORMAL}"))
    return findings


def _shape_bbox(shape) -> tuple[int, int, int, int] | None:
    """(left, top, width, height) を EMU int で返す。取得不能なら None。"""
    try:
        return (int(shape.left), int(shape.top),
                int(shape.width), int(shape.height))
    except (TypeError, ValueError):
        return None


def _shape_fill_hex(shape) -> str | None:
    """solid / gradient 塗りの主要色 hex（大文字・# なし）。無ければ None。"""
    from pptx.enum.dml import MSO_FILL_TYPE  # noqa: PLC0415
    try:
        ftype = shape.fill.type
        if ftype == MSO_FILL_TYPE.SOLID:
            return str(shape.fill.fore_color.rgb)
        if ftype == MSO_FILL_TYPE.GRADIENT:
            return str(shape.fill.gradient_stops[0].color.rgb)
    except (TypeError, AttributeError, ValueError, IndexError):
        return None
    return None


def _run_color_hex(run) -> str | None:
    """run の明示 RGB 文字色 hex（大文字）。テーマ色/継承なら None。"""
    from pptx.enum.dml import MSO_COLOR_TYPE  # noqa: PLC0415
    try:
        color = run.font.color
        if color is not None and color.type == MSO_COLOR_TYPE.RGB:
            return str(color.rgb)
    except (TypeError, AttributeError, ValueError):
        return None
    return None


def _is_overlap_candidate(shape, slide_w: int, slide_h: int, thin: int) -> bool:
    """重なり判定の対象（塗りのある通常オートシェイプ）か。
    背景・細い装飾（バー/罫線）・楕円（バッジ/ベン図）は除外する。"""
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE  # noqa: PLC0415
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            return False
        if shape.auto_shape_type == MSO_SHAPE.OVAL:
            return False
    except (TypeError, ValueError, AttributeError):
        return False
    if _shape_fill_hex(shape) is None:
        return False
    bb = _shape_bbox(shape)
    if bb is None:
        return False
    _, _, w, h = bb
    if min(w, h) < thin:                              # 細い装飾を除外
        return False
    if w >= 0.95 * slide_w and h >= 0.95 * slide_h:   # 背景を除外
        return False
    return True


def _rect_overlap_area(a: tuple, b: tuple) -> int:
    la, ta, wa, ha = a
    lb, tb, wb, hb = b
    ix = max(0, min(la + wa, lb + wb) - max(la, lb))
    iy = max(0, min(ta + ha, tb + hb) - max(ta, tb))
    return ix * iy


# ---------------------------------------------------------------------------
# PPTX チェック
# ---------------------------------------------------------------------------

def check_pptx(pptx_path: Path, label_prefix: str | None = None,
               theme: dict | None = None) -> list[Finding]:
    """生成済み .pptx を検査し、レンダリング品質問題を返す。

    theme（token→"RRGGBB"）を渡すと色トークン外の文字色も検査する。
    """
    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Inches  # noqa: PLC0415

    from deckgen import layout  # noqa: PLC0415

    prefix = label_prefix or pptx_path.as_posix()
    findings: list[Finding] = []
    prs = Presentation(str(pptx_path))

    type_scale = layout.TYPE_SCALE
    allowed_colors = {v.upper() for v in theme.values()} if theme else None
    slide_w, slide_h = int(layout.SLIDE_W), int(layout.SLIDE_H)
    eps = int(Inches(BOUNDS_EPS_IN))
    thin = int(Inches(THIN_SHAPE_IN))

    for i, slide in enumerate(prs.slides, 1):
        slide_label = f"{prefix}:slide-{i}"
        slide_sizes: set[float] = set()
        overlap_boxes: list[tuple] = []

        for shape in slide.shapes:
            # QA-PPTX-BOUNDS: スライド境界の超過（全シェイプ対象）
            bb = _shape_bbox(shape)
            if bb is not None:
                l, t, w, h = bb
                if (l < -eps or t < -eps
                        or l + w > slide_w + eps or t + h > slide_h + eps):
                    findings.append(Finding(
                        "QA-PPTX-BOUNDS", "warning", slide_label,
                        "シェイプがスライド境界を超過"))

            # 重なり判定の候補（塗り形状）を収集
            if _is_overlap_candidate(shape, slide_w, slide_h, thin):
                overlap_boxes.append(bb)

            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            text = tf.text

            # QA-PPTX-EMPTY: 空テキストフレーム
            if not text.strip():
                findings.append(Finding(
                    "QA-PPTX-EMPTY", "warning", slide_label, "空のテキストフレーム"))
                continue

            # QA-PPTX-PH: プレースホルダ残存
            if PLACEHOLDER_RE.search(text):
                findings.append(Finding(
                    "QA-PPTX-PH", "warning", slide_label, "プレースホルダ残存"))

            # QA-PPTX-OVERFLOW: テキスト量超過（面積ベースのヒューリスティック）
            try:
                w_cm = shape.width.cm
                h_cm = shape.height.cm
                if w_cm > 0 and h_cm > 0:
                    limit = max(OVERFLOW_MIN_CHARS, w_cm * h_cm * OVERFLOW_CHARS_PER_CM2)
                    if len(text) > limit:
                        findings.append(Finding(
                            "QA-PPTX-OVERFLOW", "warning", slide_label,
                            f"テキスト量超過 ({len(text)} 文字 > {int(limit)})"))
            except (AttributeError, TypeError):
                pass

            # ラン単位: フォントサイズ・色トークンの検査
            for para in tf.paragraphs:
                for run in para.runs:
                    size = run.font.size
                    if size is not None:
                        pt = round(size.pt, 1)
                        slide_sizes.add(pt)
                        # QA-PPTX-MINFONT: 最小フォント割れ
                        if pt < MIN_FONT_PT:
                            findings.append(Finding(
                                "QA-PPTX-MINFONT", "warning", slide_label,
                                f"フォント {pt}pt < {MIN_FONT_PT}pt"))
                        # QA-PPTX-TYPESCALE: 型スケール逸脱
                        elif round(pt) not in type_scale:
                            findings.append(Finding(
                                "QA-PPTX-TYPESCALE", "warning", slide_label,
                                f"型スケール外のフォント {pt}pt"))
                    # QA-PPTX-COLORTOKEN: テーマ色トークン外の文字色
                    if allowed_colors is not None:
                        chex = _run_color_hex(run)
                        if chex is not None and chex.upper() not in allowed_colors:
                            findings.append(Finding(
                                "QA-PPTX-COLORTOKEN", "warning", slide_label,
                                f"色トークン外の文字色 #{chex}"))

        # QA-PPTX-SIZEKINDS: 1 スライドのフォントサイズ種が多い
        if len(slide_sizes) > MAX_SIZE_KINDS:
            findings.append(Finding(
                "QA-PPTX-SIZEKINDS", "warning", slide_label,
                f"フォントサイズ種が多い ({len(slide_sizes)} > {MAX_SIZE_KINDS})"))

        # QA-PPTX-OVERLAP: 塗り形状同士の過剰な重なり
        for a in range(len(overlap_boxes)):
            for b in range(a + 1, len(overlap_boxes)):
                area = _rect_overlap_area(overlap_boxes[a], overlap_boxes[b])
                if area <= 0:
                    continue
                _, _, wa, ha = overlap_boxes[a]
                _, _, wb, hb = overlap_boxes[b]
                min_area = min(wa * ha, wb * hb)
                if min_area > 0 and area > OVERLAP_AREA_FRAC * min_area:
                    findings.append(Finding(
                        "QA-PPTX-OVERLAP", "warning", slide_label,
                        "塗り形状が過剰に重なっている"))
                    break
            else:
                continue
            break

    return findings


# ---------------------------------------------------------------------------
# Deck ファサード
# ---------------------------------------------------------------------------

def check_deck(slug_or_path: str) -> list[Finding]:
    """outline.yml / index.html / .pptx をまとめて検査する。"""
    findings: list[Finding] = []
    outline_path = resolve_outline_path(slug_or_path)
    deck_dir = outline_path.parent
    slug = deck_dir.name

    # outline.yml
    try:
        outline = load_outline(slug_or_path)
        findings.extend(check_outline(outline, f"{slug}/outline.yml"))
    except Exception as exc:
        findings.append(Finding(
            "QA-OUTLINE-LOAD", "error", f"{slug}/outline.yml",
            f"読み込みエラー: {exc}"))
        return findings

    # テーマを解決（色トークン検査・コントラスト検査に使う）
    from deckgen.theme import get_theme  # noqa: PLC0415
    theme_name = (outline.get("deck") or {}).get("theme")
    theme = get_theme(theme_name)

    # QA-THEME-CONTRAST: 意味色ペアの WCAG コントラスト
    findings.extend(check_theme_contrast(theme, f"{slug}/theme:{theme_name or 'default'}"))

    # index.html（存在する場合のみ）
    html_path = deck_dir / "index.html"
    if html_path.exists():
        findings.extend(check_html(html_path, label_prefix=f"{slug}/index.html"))

    # {slug}.pptx（存在する場合のみ）
    pptx_path = deck_dir / f"{slug}.pptx"
    if pptx_path.exists():
        findings.extend(check_pptx(
            pptx_path, label_prefix=f"{slug}/{slug}.pptx", theme=theme))

    return findings


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def _print_report(findings: list[Finding]) -> None:
    for f in findings:
        marker = "✗" if f.severity == "error" else "⚠"
        print(f"{marker} [{f.check}] {f.path}: {f.message}")
    errors = sum(1 for f in findings if f.severity == "error")
    warnings = sum(1 for f in findings if f.severity == "warning")
    print(f"— error: {errors} / warning: {warnings}")


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        prog="deckgen.qa",
        description="生成スライドの QA チェック（post-generation fitness function）",
    )
    parser.add_argument("target", nargs="?",
                        help="deck の slug またはパス（--all と排他）")
    parser.add_argument("--all", action="store_true",
                        help="domains/presentation/decks/ 配下の全 deck を検査")
    parser.add_argument("--json", action="store_true",
                        help="JSON で出力（check_structure.py 準拠）")
    parser.add_argument("--strict", action="store_true",
                        help="warning でも exit 1")
    args = parser.parse_args(argv)

    if args.all:
        targets = sorted(d.name for d in DECKS_DIR.iterdir() if d.is_dir())
    elif args.target:
        targets = [args.target]
    else:
        parser.print_help()
        return 2

    all_findings: list[Finding] = []
    for target in targets:
        all_findings.extend(check_deck(target))

    if args.json:
        summary = {
            "error": sum(1 for f in all_findings if f.severity == "error"),
            "warning": sum(1 for f in all_findings if f.severity == "warning"),
        }
        print(json.dumps(
            {"findings": [asdict(f) for f in all_findings], "summary": summary},
            ensure_ascii=False, indent=2))
    else:
        _print_report(all_findings)

    errors = sum(1 for f in all_findings if f.severity == "error")
    warnings = sum(1 for f in all_findings if f.severity == "warning")
    if errors:
        return 1
    if args.strict and warnings:
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
