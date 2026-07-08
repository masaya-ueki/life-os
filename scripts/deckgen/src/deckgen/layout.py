"""スライド組み立ての共通ヘルパ（寸法・色・テキストボックス・図形）。

base.css.md の設計値（16:9・本文≥24px相当・見出し・配色トークン）を
pptx のネイティブ要素に写像する。すべて編集可能な要素として生成する。
"""

from __future__ import annotations

import warnings

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from deckgen.theme import FONT

# --- スライド寸法（16:9 ワイド） ---
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# --- 型スケール（modular type scale・単一ソース） ---
# 全 expression が参照する名前付きフォントサイズ（pt）。「文字サイズの統一」要望への
# 対応として、散在していた ~15 種を 8 段に集約する。逸脱は qa.py の視覚チェックが検出する。
FONT_CAPTION = 14   # 注記・軸ラベル・日付・ステップ説明（可読性の下限）
FONT_SMALL = 16     # バッジ・表セル・孫ノード・表紙の補足行
FONT_BODY = 18      # 本文・ステップラベル・KPI ラベル・表ヘッダ
FONT_LEAD = 22      # リード(summary)・カード見出し・強調ラベル・箇条書き本文
FONT_H2 = 36        # スライド見出し（add_header のタイトル）
FONT_H1 = 40        # メッセージ・big-number の単位
FONT_DISPLAY = 48   # 表紙タイトル・KPI 数値
FONT_HERO = 96      # big-number の数値

TYPE_SCALE = frozenset(
    {FONT_CAPTION, FONT_SMALL, FONT_BODY, FONT_LEAD,
     FONT_H2, FONT_H1, FONT_DISPLAY, FONT_HERO}
)

# --- 余白グリッド（8pt ベースライン・単一ソース） ---
# 余白・ギャップ・インセットは 8pt グリッド（微細インセットのみ 4pt 半グリッド）の
# 倍数トークンで表す。「図に対する余白」要望への対応。
GRID = Pt(8)               # 基準グリッド（8pt ≈ 0.111in）
SPACE_HALF = Pt(4)         # 半グリッド（微細インセット専用）
SPACE_1 = Pt(8)            # ≈ 0.111in
SPACE_2 = Pt(16)           # ≈ 0.222in
SPACE_3 = Pt(24)           # ≈ 0.333in
SPACE_4 = Pt(32)           # ≈ 0.444in
SPACE_5 = Pt(40)           # ≈ 0.556in
SPACE_6 = Pt(48)           # ≈ 0.667in

# 角丸の一貫半径。既定の roundRect（shorter*0.16667）は大きいカードで過大になり
# 大小のカードで角丸がばらつく。サイズ非依存の一定半径に揃えて整合させる。
CARD_RADIUS = Inches(0.14)

# --- カード統一スタイル（単一の真実） ---
# 「カード」風の面（comparison カラム / KPI / flow ステップ / tree ノード /
# matrix 象限）はすべて add_card() で描く。expression ごとに枠線色・太さ・
# 角丸有無がばらつくのを防ぐため、スタイル差は variant / accent_bar だけで表す。
CARD_LINE_W = 1.0            # カード枠線太さ（全カード共通・pt 数値。Pt() は適用側で行う）
CARD_BAR_W = Inches(0.07)    # 左端アクセントバー幅
CARD_PAD_X = Pt(12)          # カード内テキストの左右余白（1.5 グリッド）
CARD_PAD_Y = Pt(12)          # カード内テキストの上下余白（1.5 グリッド）
CARD_TEXT_LEFT = CARD_BAR_W + CARD_PAD_X  # アクセントバー付きカードの文字開始位置
CARD_LABEL_H = SPACE_5       # カード見出し行の高さ（40pt）

# --- 余白とリージョン（8pt グリッド整列） ---
MARGIN = Pt(44)            # 5.5×8pt（≈0.61in）
CONTENT_LEFT = MARGIN
CONTENT_WIDTH = SLIDE_W - 2 * MARGIN

TITLE_TOP = Pt(32)         # 4×8pt（≈0.44in）
TITLE_HEIGHT = Pt(64)      # 8×8pt（≈0.89in）
LEAD_TOP = Pt(104)         # 13×8pt（≈1.44in）
LEAD_HEIGHT = Pt(56)       # 7×8pt（≈0.78in）
BODY_TOP = Pt(168)         # 21×8pt（≈2.33in）
BODY_BOTTOM_MARGIN = Pt(32)  # 4×8pt（≈0.44in）
BODY_HEIGHT = SLIDE_H - BODY_TOP - BODY_BOTTOM_MARGIN

# Region = (left, top, width, height)（すべて EMU int）
Region = tuple

# --- 図解共通定数（flow / structure / comparison が参照する共有値） ---
# 枠線太さはカード統一スタイルの CARD_LINE_W に一本化（add_card / add_box_shape 既定）。

# ボックス内余白（8pt グリッド整列）
DIAG_PAD_XS = SPACE_HALF    # 4pt: 極小余白（矢印-ボックス間）
DIAG_PAD_SM = SPACE_1       # 8pt: 小余白（バッジオフセット）
DIAG_PAD = Pt(12)           # 12pt(1.5グリッド): 標準余白（テキスト左右）

# ステップバッジ（flow の番号円）
BADGE_D = Inches(0.45)      # バッジ直径
BADGE_FONT = FONT_SMALL     # バッジ数字フォントサイズ

# フロー矢印
FLOW_ARROW_W = Inches(0.50)  # 水平矢印幅
FLOW_ARROW_H = Inches(0.35)  # 垂直矢印高さ
FLOW_LABEL_FONT = FONT_BODY  # ステップラベルフォントサイズ
FLOW_DESC_FONT = FONT_CAPTION  # ステップ説明フォントサイズ
FLOW_H_GAP_PAD = SPACE_1     # 8pt: gap の内訳（FLOW_ARROW_W + これ = 合計 gap）
FLOW_LABEL_PAD = SPACE_1     # 8pt: ラベルテキストボックス左右余白
FLOW_ARROW_OFFSET_H = SPACE_HALF  # 4pt: ボックス→矢印 x オフセット
FLOW_ARROW_OFFSET_V = Pt(2)  # 垂直フロー: ボックス→矢印 y オフセット（微調整）
BADGE_Y_OFFSET = SPACE_1     # 8pt: バッジ y オフセット

# ツリー図
TREE_ROOT_W = Inches(4.2)    # ルートボックス最大幅
TREE_ROOT_H = Inches(0.75)   # ルートボックス高さ
TREE_CHILD_W = Inches(3.4)   # 子ボックス最大幅
TREE_CHILD_MIN_H = Inches(0.7)  # 子ボックス最小高さ
TREE_NODE_GAP = SPACE_2      # 16pt: 子ノード間ギャップ
TREE_VERT_SPAN = Inches(0.85) # ルート下端→子ボックス上端の距離
TREE_BUS_OFFSET = Inches(0.42) # バスラインのルート下端からのオフセット
TREE_ROOT_FONT = FONT_BODY   # ルートノードフォントサイズ
TREE_CHILD_FONT = FONT_SMALL  # 子ノードフォントサイズ
TREE_GC_FONT = FONT_CAPTION  # 孫ノードフォントサイズ

# フローの内部レイアウト
FLOW_BOX_MAX_H = Inches(2.4)   # 水平ステップボックスの最大高さ
FLOW_LABEL_TOP = BADGE_Y_OFFSET + BADGE_D + DIAG_PAD_XS  # バッジ下端 + 4pt
FLOW_LABEL_H = Inches(0.65)    # ステップラベル行の高さ
FLOW_ARROW_BODY_H = Inches(0.4)  # 水平矢印の高さ / 垂直矢印の幅
FLOW_V_TEXT_LEFT = Inches(0.9)   # 垂直フロー: バッジと重ならない文字開始位置

# マトリクス（2x2）
MATRIX_AXIS_PAD_L = Inches(0.5)  # Y 軸ラベル用の左マージン
MATRIX_AXIS_PAD_B = Inches(0.4)  # X 軸ラベル用の下マージン
MATRIX_CELL_INSET = SPACE_HALF   # 4pt: 象限セル同士の間隔（片側）

# 強調（emphasis）・KPI
EMPHASIS_PAD_X = SPACE_6         # 48pt: アクセント面内テキストの左右余白
KPI_CARD_MAX_H = Inches(2.6)     # KPI カードの最大高さ

# 注記（comparison / chart の note 共通）
NOTE_H = Inches(0.45)            # 注記行の高さ
NOTE_GAP = SPACE_1               # 本文と注記の間隔

# 比較表（comparison mode=table）
TABLE_FIRST_COL_W = Inches(2.6)  # 評価軸（1列目）の幅
TABLE_ROW_MAX_H = Inches(0.7)    # 行の最大高さ

# --- 表紙（title）--- コンテンツスライドと左揃えラインを共有する
COVER_BAR_W = Inches(0.18)       # 左端アクセントバー幅
COVER_TITLE_TOP = Inches(2.0)    # タイトルボックス上端
COVER_TITLE_H = Inches(2.0)      # タイトルボックス高さ（下端 = ルール線）
COVER_RULE_Y = COVER_TITLE_TOP + COVER_TITLE_H
COVER_SUBTITLE_TOP = COVER_RULE_Y + SPACE_1
COVER_SUBTITLE_H = Inches(1.0)
COVER_CONTENT_TOP = COVER_SUBTITLE_TOP + COVER_SUBTITLE_H + SPACE_1
COVER_CONTENT_H = Inches(1.5)
COVER_DATE_H = NOTE_H
COVER_DATE_TOP = SLIDE_H - COVER_DATE_H - SPACE_3


def rgb(hexcolor: str) -> RGBColor:
    return RGBColor.from_string(hexcolor)


def blank_layout(prs):
    """テンプレ非依存で使える Blank に近いレイアウトを返す。"""
    for layout in prs.slide_layouts:
        if layout.name.strip().lower() == "blank":
            return layout
    # 慣例上 index 6 が Blank。無ければ最後を使う。
    layouts = list(prs.slide_layouts)
    if len(layouts) > 6:
        return layouts[6]
    return layouts[-1]


def add_slide(prs):
    return prs.slides.add_slide(blank_layout(prs))


def fill_background(slide, hexcolor: str) -> None:
    """スライド全面に背景矩形を敷く（テンプレ無し時のテーマ背景）。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(hexcolor)
    shape.line.fill.background()
    shape.shadow.inherit = False
    # 背景は最背面へ
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    *,
    size,
    color,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=FONT,
    wrap=True,
    autofit=False,
):
    """テキストボックス。size / color は必須（テーマトークンを明示的に渡す）。
    autofit=True で溢れ時に PowerPoint がテキストを自動縮小する（normAutofit）。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if autofit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    _style_run(p.add_run(), text, size, color, bold, font)
    return box


def _style_run(run, text, size, color, bold, font):
    run.text = text
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.name = font
    f.color.rgb = rgb(color)


def add_run(paragraph, text, *, size, color, bold=False, font=FONT):
    """段落へスタイル済み run を追加する（expression 側の重複実装を排除する公開 API）。"""
    run = paragraph.add_run()
    _style_run(run, text, size, color, bold, font)
    return run


def set_center_text(shape, text, *, size, color, bold=True, font=FONT):
    """図形の text_frame に中央寄せ（水平・垂直）のテキストを設定し、
    追記用に text_frame を返す。カード・バッジ・ピラミッド層などで共用する。"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, text, size=size, color=color, bold=bold, font=font)
    return tf


def _set_hanging_indent(paragraph, mar_left: int, indent: int) -> None:
    """段落にぶら下げインデントを設定する（marL / indent を EMU で指定）。
    折り返し行を bullet ではなくテキスト先頭に揃える。indent は負値。"""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(int(mar_left)))
    pPr.set("indent", str(int(indent)))


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items,
    *,
    size,
    color,
    font=FONT,
    bullet="•  ",
    line_spacing=1.3,
    space_after=8,
    autofit=False,
):
    """箇条書きテキストフレーム。items は str か (text, level) のタプル列。
    size / color は必須（テーマトークンを明示的に渡す）。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    if autofit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        prefix = bullet if bullet else ""
        # ぶら下げインデント: 折り返し行を bullet/マーク幅ぶん右へ揃える
        hang = int(Pt(size * (1.1 if bullet else 0.9)))
        level_indent = level * int(Pt(size * 1.4))
        _set_hanging_indent(p, level_indent + hang, -hang)
        _style_run(p.add_run(), f"{prefix}{text}", size, color, False, font)
    return box


def add_box_shape(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill,
    line=None,
    line_width=CARD_LINE_W,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE,
):
    """低レベルの図形描画。fill は必須（テーマトークンを明示的に渡す）。
    「カード」風の面は原則 add_card() を使うこと（スタイル統一のため）。"""
    sp = slide.shapes.add_shape(shape, left, top, width, height)
    # 角丸矩形はサイズ非依存の一貫半径に揃える（既定は shorter*0.16667 で過大）
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            ss = min(int(width), int(height))
            if ss > 0:
                sp.adjustments[0] = min(0.5, int(CARD_RADIUS) / ss)
        except (TypeError, ValueError, IndexError):
            pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = rgb(fill)
    if line:
        sp.line.color.rgb = rgb(line)
        sp.line.width = Pt(line_width)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_card(slide, left, top, width, height, theme, *, variant="outline",
             accent_bar=None):
    """統一スタイルの「カード」を描く（角丸・枠線・アクセントバーの単一の真実）。

    variant:
      - "outline": card 塗り + line 枠 CARD_LINE_W の標準カード（既定）
      - "accent" : accent 塗り・枠なし。スライド内で最重要の面にだけ使う
    accent_bar: 色 hex を渡すと左端にピル形状のバーを添える（カラム見出し色など）。
    """
    if variant not in ("outline", "accent"):
        raise ValueError(f"add_card: 未知の variant '{variant}'（outline / accent のみ）")
    if variant == "accent":
        sp = add_box_shape(slide, left, top, width, height,
                           fill=theme["accent"], line=None)
    else:
        sp = add_box_shape(slide, left, top, width, height,
                           fill=theme["card"], line=theme["line"],
                           line_width=CARD_LINE_W)
    if accent_bar:
        add_accent_bar(slide, left, top, height, accent_bar,
                       inset=int(CARD_RADIUS))
    return sp


def add_freeform_polygon(slide, points, *, fill, line=None, line_width=CARD_LINE_W):
    """EMU 座標の点列 [(x, y), ...] から閉じた多角形（自由形状）を描く。
    ピラミッドの連続三角形シルエットなど、矩形以外の図形に使う。"""
    x0, y0 = points[0]
    fb = slide.shapes.build_freeform(int(x0), int(y0), scale=1.0)
    fb.add_line_segments([(int(x), int(y)) for x, y in points[1:]], close=True)
    sp = fb.convert_to_shape()
    sp.fill.solid()
    sp.fill.fore_color.rgb = rgb(fill)
    if line:
        sp.line.color.rgb = rgb(line)
        sp.line.width = Pt(line_width)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_arrow(slide, left, top, width, height, color):
    """右向き矢印。color は必須（通常 theme["muted"]）。"""
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = rgb(color)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_down_arrow(slide, left, top, width, height, color):
    """下向き矢印。color は必須（通常 theme["muted"]）。"""
    sp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = rgb(color)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_rule(slide, left, top, width, color, weight=2.5):
    """水平の罫線（タイトル下線など）。"""
    return add_connector(slide, left, top, left + width, top, color, weight)


def add_connector(slide, x1, y1, x2, y2, color, weight=1.5):
    """任意の2点を結ぶ直線コネクタ（ツリーの枝・見出し下線など）。"""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(weight)
    # 既定のテーマ影を無効化（下線・枝に不要なドロップシャドウが乗るのを防ぐ）
    conn.shadow.inherit = False
    return conn


def add_gradient_fill(shape, color1: str, color2: str, angle: float = 135.0) -> None:
    """ソリッド塗りをリニアグラデーションに変換する（python-pptx 1.0+）。
    color1 が始点、color2 が終点。angle=135 で左上→右下方向。
    """
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    stops = fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = rgb(color1)
    stops[1].position = 1.0
    stops[1].color.rgb = rgb(color2)


def add_accent_bar(slide, left, top, height, color: str, width: float = None,
                   inset=0) -> None:
    """カード左端に細いアクセントバーを追加してデプス感を出す。
    inset>0 でカードの角丸半径ぶん上下を詰め、丸角からのはみ出しを防ぐ。
    端は丸めたピル形状にしてカードの角丸と整合させる。"""
    bar_w = width if width is not None else Inches(0.07)
    y = top + inset
    h = height - 2 * inset
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, bar_w, h)
    try:
        sp.adjustments[0] = 0.5  # 半径 = 幅の半分 → 端が半円のピル形状
    except (IndexError, ValueError):
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = rgb(color)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def set_fill_alpha(shape, pct: float) -> None:
    """ソリッド塗りの不透明度を pct%（0–100）に設定する。

    python-pptx は塗りの透過を直接サポートしないため、`a:srgbClr` に
    `a:alpha` 子要素を付与する。ベン図の重なり表現などに使う。
    """
    spPr = shape._element.spPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is None:
        warnings.warn("set_fill_alpha: solidFill が無いため透過を適用できません", stacklevel=2)
        return
    srgb = solidFill.find(qn("a:srgbClr"))
    if srgb is None:
        warnings.warn("set_fill_alpha: srgbClr が無いため透過を適用できません", stacklevel=2)
        return
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))


def style_cell(cell, text, *, size=FONT_BODY, color, bold=False, fill=None,
               align=PP_ALIGN.LEFT, font=FONT):
    """表セルのスタイル。color は必須（テーマトークンを明示的に渡す）。"""
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(fill)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.1)
    cell.margin_right = Inches(0.1)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _style_run(p.add_run(), text, size, color, bold, font)


def add_table(slide, left, top, width, height, rows, cols):
    gf = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = gf.table
    # python-pptx 既定のテーマ縞模様を抑える（テーマ色で塗るため）
    table.first_row = False
    table.horz_banding = False
    return table


def add_header(slide, theme, title, summary):
    """各コンテンツスライド共通の見出し（タイトル + 下線 + リード）を描き、
    本文用 Region (left, top, width, height) を返す。"""
    add_textbox(
        slide,
        CONTENT_LEFT,
        TITLE_TOP,
        CONTENT_WIDTH,
        TITLE_HEIGHT,
        title,
        size=FONT_H2,
        color=theme["accent"],
        bold=True,
        anchor=MSO_ANCHOR.BOTTOM,
    )
    add_rule(
        slide,
        CONTENT_LEFT,
        TITLE_TOP + TITLE_HEIGHT,
        CONTENT_WIDTH,
        theme["line"],
        weight=2.5,
    )
    if summary:
        add_textbox(
            slide,
            CONTENT_LEFT,
            LEAD_TOP,
            CONTENT_WIDTH,
            LEAD_HEIGHT,
            summary,
            size=FONT_LEAD,
            color=theme["fg"],
            bold=True,
        )
    return (CONTENT_LEFT, BODY_TOP, CONTENT_WIDTH, BODY_HEIGHT)
