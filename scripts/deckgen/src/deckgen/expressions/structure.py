"""structure: 関係・分類・位置づけ。
type = matrix-2x2 | tree | pyramid | matrix-table | venn。
matrix-table は rows/cols を持たない汎用分類のため tree 描画にフォールバックする。

data 契約: slide-expression/references/structure.md
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from deckgen import layout


def render(pslide, theme, slide, region):
    data = slide.get("data") or {}
    t = data.get("type", "tree")
    if t == "matrix-2x2":
        _matrix_2x2(pslide, theme, data, region)
    elif t == "pyramid":
        _pyramid(pslide, theme, data, region)
    elif t == "venn":
        _venn(pslide, theme, data, region)
    elif t == "matrix-table":
        _matrix_table(pslide, theme, data, region)
    else:  # tree / その他
        _tree(pslide, theme, data, region)


def _centered_text(shape, text, size, color, *, bold=True):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = layout.FONT
    run.font.color.rgb = layout.rgb(color)


def _matrix_2x2(pslide, theme, data, region):
    left, top, width, height = region
    axis_x = str(data.get("axis_x", ""))
    axis_y = str(data.get("axis_y", ""))
    quad = [str(q) for q in (data.get("quadrants") or [])]
    quad += [""] * (4 - len(quad))

    # 軸ラベル用に外周マージンを確保
    pad_l = Inches(0.5)
    pad_b = Inches(0.4)
    grid_left = left + pad_l
    grid_w = width - pad_l
    grid_h = height - pad_b
    cell_w = grid_w // 2
    cell_h = grid_h // 2
    # 契約(structure.md): quadrants は [右上, 左上, 右下, 左下] の順。
    # 右上(axis_x/axis_y がともに高い)= 最優先を accent で強調する。
    #   quadrants[0] → 右上(0,1, accent)  quadrants[1] → 左上(0,0)
    #   quadrants[2] → 右下(1,1)          quadrants[3] → 左下(1,0)
    positions = [(0, 1), (0, 0), (1, 1), (1, 0)]
    for idx, (rr, cc) in enumerate(positions):
        x = grid_left + cc * cell_w
        y = top + rr * cell_h
        top_right = rr == 0 and cc == 1
        fill = theme["accent"] if top_right else theme["card"]
        color = theme["on_accent"] if top_right else theme["fg"]
        box = layout.add_box_shape(
            pslide, x + Inches(0.05), y + Inches(0.05),
            cell_w - Inches(0.1), cell_h - Inches(0.1),
            fill=fill, line=theme["line"], line_width=1.0,
            shape=MSO_SHAPE.RECTANGLE,
        )
        _centered_text(box, quad[idx], layout.FONT_SMALL, color)
    # 軸ラベル: X は下中央(→ 右ほど高い)、Y は左縦(↑ 上ほど高い)
    if axis_x:
        layout.add_textbox(
            pslide, grid_left, top + grid_h + Inches(0.02), grid_w, pad_b,
            f"→ {axis_x}", size=layout.FONT_CAPTION, color=theme["muted"], bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
    if axis_y:
        layout.add_textbox(
            pslide, left - Inches(0.05), top, pad_l + Inches(0.05), grid_h,
            f"↑ {axis_y}", size=layout.FONT_CAPTION, color=theme["muted"], bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=True,
        )


def _tree(pslide, theme, data, region):
    root = data.get("root")
    children = data.get("children") or []
    # root + 第1階層がそろうときはコネクタ線付きのネイティブツリーで描く。
    if root and children:
        _tree_native(pslide, theme, str(root), children, region)
        return
    _tree_bullets(pslide, theme, root, children, region)


def _tree_native(pslide, theme, root, children, region):
    left, top, width, height = region
    n = len(children)

    # root ボックス(上部中央)
    root_w = min(width, layout.TREE_ROOT_W)
    root_h = layout.TREE_ROOT_H
    root_x = left + (width - root_w) // 2
    root_box = layout.add_box_shape(
        pslide, root_x, top, root_w, root_h,
        fill=theme["accent"], line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    _centered_text(root_box, root, layout.TREE_ROOT_FONT, theme["on_accent"])

    # 子ボックスを下段に横並び
    gap = layout.TREE_NODE_GAP
    child_top = top + root_h + layout.TREE_VERT_SPAN
    child_h = max(Inches(0.7), top + height - child_top)
    cell_w = (width - gap * (n - 1)) // n if n else width
    child_w = min(cell_w, layout.TREE_CHILD_W)
    root_cx = root_x + root_w // 2
    root_bottom = top + root_h
    bus_y = root_bottom + layout.TREE_BUS_OFFSET

    centers = []
    for i, node in enumerate(children):
        cell_x = left + i * (cell_w + gap)
        cx = cell_x + cell_w // 2
        centers.append(cx)
        bx = cx - child_w // 2
        box = layout.add_box_shape(
            pslide, bx, child_top, child_w, child_h,
            fill=theme["card"], line=theme["line"], line_width=1.0,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        _fill_child(box, theme, node)
        # 枝: バス線から各子ボックス上端へ
        layout.add_connector(pslide, cx, bus_y, cx, child_top, theme["muted"], 1.5)

    # 枝: root 下端 → バス線(縦) と 子をつなぐ横バス
    layout.add_connector(pslide, root_cx, root_bottom, root_cx, bus_y, theme["muted"], 1.5)
    if centers:
        layout.add_connector(pslide, centers[0], bus_y, centers[-1], bus_y, theme["muted"], 1.5)


def _fill_child(box, theme, node):
    """子ボックスに名前(太字)と、孫があれば小さな一覧を書き込む。"""
    if isinstance(node, dict):
        name = str(node.get("name", ""))
        grandchildren = [
            str(g.get("name", "") if isinstance(g, dict) else g)
            for g in (node.get("children") or [])
        ]
    else:
        name, grandchildren = str(node), []
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = name
    r.font.size = Pt(layout.TREE_CHILD_FONT)
    r.font.bold = True
    r.font.name = layout.FONT
    r.font.color.rgb = layout.rgb(theme["fg"])
    for g in grandchildren:
        gp = tf.add_paragraph()
        gp.alignment = PP_ALIGN.CENTER
        gr = gp.add_run()
        gr.text = f"・{g}"
        gr.font.size = Pt(layout.TREE_GC_FONT)
        gr.font.name = layout.FONT
        gr.font.color.rgb = layout.rgb(theme["muted"])


def _tree_bullets(pslide, theme, root, children, region):
    left, top, width, height = region
    items = []
    if root:
        items.append((str(root), 0))
    _walk(children or [], 1, items)
    if not items:
        return
    layout.add_bullets(
        pslide, left, top, width, height, items,
        size=layout.FONT_LEAD, color=theme["fg"], bullet="", line_spacing=1.3, space_after=8,
    )


def _walk(nodes, level, out):
    for node in nodes:
        if isinstance(node, dict):
            name = node.get("name", "")
            mark = "▸ " if level == 1 else "・"
            out.append((f"{mark}{name}", level))
            _walk(node.get("children") or [], level + 1, out)
        else:
            out.append((f"・{node}", level))


def _venn(pslide, theme, data, region):
    left, top, width, height = region
    sets = [str(s) for s in (data.get("sets") or [])]
    overlap = str(data.get("overlap", ""))
    # 2集合の重なり円を描く。集合名が無ければ tree フォールバック。
    if len(sets) < 2:
        _tree(pslide, theme, data, region)
        return

    overlap_frac = 0.34
    diam = min(int(height) - int(Inches(0.2)), int(width * 0.58))
    total_w = int(diam * (2 - overlap_frac))
    start_x = left + (width - total_w) // 2
    cy = top + (height - diam) // 2
    x1 = start_x
    x2 = start_x + int(diam * (1 - overlap_frac))

    for cx, fill in ((x1, theme["accent"]), (x2, theme["accent2"])):
        circle = layout.add_box_shape(
            pslide, cx, cy, diam, diam,
            fill=fill, line=None, shape=MSO_SHAPE.OVAL,
        )
        layout.set_fill_alpha(circle, 55)

    label_h = Inches(0.6)
    label_cy = cy + diam // 2 - label_h // 2
    # 左集合ラベル(重なりの外・左寄り)
    layout.add_textbox(
        pslide, x1 - Inches(0.1), label_cy, int(diam * (1 - overlap_frac) * 0.8), label_h,
        sets[0], size=layout.FONT_SMALL, color=theme["on_accent"], bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # 右集合ラベル(重なりの外・右寄り)
    right_only_x = x2 + int(diam * overlap_frac)
    layout.add_textbox(
        pslide, right_only_x, label_cy, int(diam * (1 - overlap_frac) * 0.8), label_h,
        sets[1], size=layout.FONT_SMALL, color=theme["on_accent"], bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # 重なりラベル(中央)
    if overlap:
        ov_x = x2
        ov_w = (x1 + diam) - x2
        layout.add_textbox(
            pslide, ov_x, label_cy, ov_w, label_h,
            overlap, size=layout.FONT_CAPTION, color=theme["fg"], bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )


def _color_luminance(hex_color: str) -> float:
    """0–255 の相対輝度を返す（高いほど明るい）。"""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return 0.299 * r + 0.587 * g + 0.114 * b


def _pyramid(pslide, theme, data, region):
    left, top, width, height = region
    layers = [str(x) for x in (data.get("layers") or [])]
    if not layers:
        return
    n = len(layers)
    gap = layout.SPACE_1
    row_h = (height - gap * (n - 1)) // n
    apex_x = left + width // 2      # 頂点の x（上中央）
    span = height                   # 頂点(top)→底辺(top+height) の高さ

    def half_w_at(y):
        # 頂点(y=top)で 0、底辺(y=top+height)で width/2 の連続直線上の半幅。
        frac = (y - top) / span if span else 0
        return int(width // 2 * frac)

    shade = [theme["accent"], theme["accent2"], theme["muted"], theme["line"]]
    # layers[0]=土台 を最下段に。reversed で i=0 が頂点(最上段・最狭)。
    for i, layer in enumerate(reversed(layers)):
        y_top = top + i * (row_h + gap)
        y_bot = y_top + row_h
        htw_top = half_w_at(y_top)   # 上辺の半幅（各層の端が頂点‐底辺の2直線に載る）
        htw_bot = half_w_at(y_bot)   # 下辺の半幅
        # 上辺→下辺の台形（頂点段は上辺幅0=三角形）。連続した三角形シルエットになる。
        points = [
            (apex_x - htw_top, y_top),
            (apex_x + htw_top, y_top),
            (apex_x + htw_bot, y_bot),
            (apex_x - htw_bot, y_bot),
        ]
        fill = shade[i % len(shade)]
        box = layout.add_freeform_polygon(pslide, points, fill=fill)
        # 塗り色の輝度で文字色を切り替え（薄い背景には濃い文字）
        text_color = theme["fg"] if _color_luminance(fill) > 160 else theme["on_accent"]
        _centered_text(box, layer, layout.FONT_SMALL, text_color)


def _matrix_table(pslide, theme, data, region):
    # rows/cols が無い汎用分類表は tree フォールバックで十分
    _tree(pslide, theme, data, region)
