"""emphasis: 1枚1メッセージの強調。
mode = big-number | kpi | message | quote。

data 契約: slide-expression/references/emphasis.md
"""

from __future__ import annotations

import warnings

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from deckgen import layout


def render(pslide, theme, slide, region):
    data = slide.get("data") or {}
    mode = data.get("mode", "message")
    if mode == "big-number":
        _big_number(pslide, theme, data, region)
    elif mode == "kpi":
        _kpi(pslide, theme, data, region)
    elif mode == "quote":
        _quote(pslide, theme, data, region, slide)
    else:
        _message(pslide, theme, data, region, slide)


def _accent_card(pslide, theme, region):
    """アクセント色のグラデーションカード（最重要メッセージの面）。"""
    left, top, width, height = region
    card = layout.add_card(pslide, left, top, width, height, theme,
                           variant="accent")
    # accent → accent2 のグラデーションで奥行き感を演出
    accent2 = theme.get("accent2", theme["accent"])
    try:
        layout.add_gradient_fill(card, theme["accent"], accent2, angle=135.0)
    except Exception as exc:  # noqa: BLE001 — solid fill フォールバックの意図を優先
        # グラデーション未対応環境（旧 python-pptx 等）では solid fill のまま進める。
        # フォールバック自体は維持しつつ、失敗を warnings で可視化する（沈黙させない）。
        warnings.warn(f"emphasis: グラデーション適用に失敗（solid で継続）: {exc}",
                      stacklevel=2)
    return card


def _center_paragraph(pslide, left, top, width, height):
    """中央寄せ段落を持つテキストボックスを作り、(text_frame, 段落) を返す。"""
    box = pslide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    return tf, p


def _big_number(pslide, theme, data, region):
    left, top, width, height = region
    _accent_card(pslide, theme, region)
    value = str(data.get("value", ""))
    unit = str(data.get("unit", ""))
    label = str(data.get("label", ""))
    tf, p = _center_paragraph(pslide, left, top, width, height)
    layout.add_run(p, value, size=layout.FONT_HERO, color=theme["on_accent"],
                   bold=True)
    if unit:
        layout.add_run(p, unit, size=layout.FONT_H1, color=theme["on_accent"],
                       bold=True)
    if label:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        layout.add_run(p2, label, size=layout.FONT_LEAD,
                       color=theme["on_accent"])


def _message(pslide, theme, data, region, slide):
    left, top, width, height = region
    _accent_card(pslide, theme, region)
    text = str(data.get("text") or slide.get("summary") or "")
    _tf, p = _center_paragraph(
        pslide, left + layout.EMPHASIS_PAD_X, top,
        width - 2 * layout.EMPHASIS_PAD_X, height,
    )
    layout.add_run(p, text, size=layout.FONT_H1, color=theme["on_accent"],
                   bold=True)


def _quote(pslide, theme, data, region, slide):
    left, top, width, height = region
    _accent_card(pslide, theme, region)
    text = str(data.get("text") or slide.get("summary") or "")
    cite = str(data.get("cite", ""))
    tf, p = _center_paragraph(
        pslide, left + layout.EMPHASIS_PAD_X, top,
        width - 2 * layout.EMPHASIS_PAD_X, height,
    )
    layout.add_run(p, f"“{text}”", size=layout.FONT_H2,
                   color=theme["on_accent"], bold=True)
    if cite:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        layout.add_run(p2, cite, size=layout.FONT_LEAD,
                       color=theme["on_accent"])


def _delta_color(theme, delta: str) -> str:
    """増減表示の意味色。マイナス方向は bad、それ以外は good。"""
    if delta.lstrip().startswith(("-", "−", "▼", "▽", "↓")):
        return theme["bad"]
    return theme["good"]


def _kpi(pslide, theme, data, region):
    left, top, width, height = region
    cards = data.get("cards") or []
    if not cards:
        return
    n = len(cards)
    gap = layout.SPACE_3
    card_w = (width - gap * (n - 1)) // n
    card_h = min(layout.KPI_CARD_MAX_H, height)
    card_top = top + (height - card_h) // 2
    for i, card in enumerate(cards):
        x = left + i * (card_w + gap)
        box = layout.add_card(pslide, x, card_top, card_w, card_h, theme)
        tf = layout.set_center_text(
            box, str(card.get("num", "")),
            size=layout.FONT_DISPLAY, color=theme["accent"],
        )
        delta = str(card.get("delta", ""))
        if delta:
            layout.add_run(tf.paragraphs[0], f" {delta}",
                           size=layout.FONT_LEAD,
                           color=_delta_color(theme, delta), bold=True)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        layout.add_run(p2, str(card.get("label", "")),
                       size=layout.FONT_BODY, color=theme["fg"])
